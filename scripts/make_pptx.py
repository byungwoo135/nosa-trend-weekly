#!/usr/bin/env python3
"""주간 노사동향 리포트 JSON을 편집 가능한 PowerPoint(.pptx)로 변환한다.

Usage: .venv/bin/python scripts/make_pptx.py reports/{SLUG}.json reports/{SLUG}.pptx
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

FONT = "Malgun Gothic"

PURPLE = RGBColor(0x5F, 0x00, 0x80)
PURPLE_2 = RGBColor(0x3D, 0x00, 0x53)
ORANGE = RGBColor(0xFF, 0x7B, 0x3C)
ORANGE_TEXT = RGBColor(0xB8, 0x50, 0x1A)
INK = RGBColor(0x16, 0x21, 0x3E)
SUB = RGBColor(0x5B, 0x64, 0x74)
LINE = RGBColor(0xE7, 0xE9, 0xEE)
PAPER = RGBColor(0xF7, 0xF8, 0xFB)
RETAIL_BG = RGBColor(0xF2, 0xE6, 0xF7)
OTHER_BG = RGBColor(0xFD, 0xEE, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOURCE_GRAY = RGBColor(0x97, 0xA0, 0xB3)
AVATAR_BG = RGBColor(0xEA, 0xE3, 0xF0)
DASH_BORDER = RGBColor(0xC9, 0xB3, 0xD8)

W = 7.5  # inches, slide width


def no_shadow(shape):
    shape.shadow.inherit = False


def add_rect(slide, x, y, w, h, fill, line=None, rounded=False, radius=0.08):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = radius
        except (IndexError, ValueError):
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    no_shadow(shp)
    return shp


def add_gradient_rect(slide, x, y, w, h, c1, c2, angle=135):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.line.fill.background()
    no_shadow(shp)
    shp.fill.gradient()
    stops = shp.fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    shp.fill.gradient_angle = angle
    return shp


def set_font(run, size, color, bold=False, italic=False, font=FONT):
    """python-pptx's run.font.name only sets the *Latin* typeface (<a:latin>).
    Korean text is rendered via the East Asian typeface (<a:ea>), which falls
    back to the theme default if not set explicitly -- that mismatch is what
    causes PowerPoint to substitute a different font (and different glyph
    spacing/kerning) for the Korean text even though run.font.name looks
    right. Set latin/ea/cs to the same face so PowerPoint doesn't substitute
    anything, on any machine that has this font."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font  # sets <a:latin>
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, italic=False, font=FONT, hyperlink=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold=bold, italic=italic, font=font)
    if hyperlink:
        run.hyperlink.address = hyperlink
    return tb


def add_circular_picture(slide, image_path, x, y, d):
    """Insert a picture cropped to a circle. Falls back to a plain square
    picture if the OOXML geometry swap fails for any reason."""
    pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), Inches(d), Inches(d))
    try:
        spPr = pic._element.spPr
        old_geom = spPr.find(qn('a:prstGeom'))
        if old_geom is not None:
            spPr.remove(old_geom)
        geom = spPr.makeelement(qn('a:prstGeom'), {'prst': 'ellipse'})
        av_lst = geom.makeelement(qn('a:avLst'), {})
        geom.append(av_lst)
        spPr.append(geom)
    except Exception:
        pass
    return pic


def add_pill(slide, x, y, w, h, text, fill, text_color, size=10):
    shp = add_rect(slide, x, y, w, h, fill, rounded=True, radius=0.5)
    tf = shp.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, size, text_color, bold=True)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shp


def char_width_in(size_pt, bold=False):
    """Rough worst-case width of one Korean (full-width) character at size_pt,
    inflated a bit extra since we cannot visually verify wrapping in this
    environment (no LibreOffice) -- better to over-estimate line count than
    risk clipped/overlapping text."""
    factor = 1.25 if bold else 1.15
    return (size_pt / 72.0) * factor


def lines_needed(text, width_in, size_pt, bold=False):
    if not text:
        return 1
    cpl = max(1, int(width_in / char_width_in(size_pt, bold)))
    return max(1, -(-len(text) // cpl))


HEADLINE_SIZE = 12.5
SUMMARY_SIZE = 10
HEADLINE_LINE_H = 0.26
SUMMARY_LINE_H = 0.19
BADGE_W = 0.62
BADGE_GAP = 0.08


def build_section(slide, x, y, w, title, tag_text, items, accent, accent_bg, dot_color):
    pad = 0.22
    inner_x = x + pad
    inner_w = w - pad * 2
    cursor = y + pad

    # section title + tag
    add_text(slide, inner_x, cursor, inner_w - 1.1, 0.32, title, 16, INK, bold=True)
    add_pill(slide, x + w - pad - 1.05, cursor + 0.02, 1.05, 0.28, tag_text, accent_bg, accent, size=10)
    cursor += 0.46

    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(inner_x), Inches(cursor + 0.06), Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        no_shadow(dot)

        text_x = inner_x + 0.22
        text_w = inner_w - 0.22

        has_badge = bool(item.get("kurly_related"))
        headline_w = text_w - (BADGE_W + BADGE_GAP) if has_badge else text_w
        headline_lines = lines_needed(item["headline"], headline_w, HEADLINE_SIZE, bold=True)
        headline_h = HEADLINE_LINE_H * headline_lines
        if has_badge:
            add_pill(slide, text_x + text_w - BADGE_W, cursor, BADGE_W, 0.22, "CHECK", PURPLE, WHITE, size=8.5)
        add_text(slide, text_x, cursor, headline_w, headline_h, item["headline"], HEADLINE_SIZE, INK, bold=True,
                  hyperlink=item.get("source_url"))
        cursor += headline_h + 0.05

        summary_lines = lines_needed(item["summary"], text_w, SUMMARY_SIZE)
        summary_h = SUMMARY_LINE_H * summary_lines
        add_text(slide, text_x, cursor, text_w, summary_h, item["summary"], SUMMARY_SIZE, SUB)
        cursor += summary_h + 0.04

        add_text(slide, text_x, cursor, text_w, 0.18, f'{item["source_name"]} · {item["published"]}', 8.5, SOURCE_GRAY)
        cursor += 0.18 + 0.16  # gap before next item

    cursor += pad - 0.14
    box_h = cursor - y
    return box_h


def main():
    if len(sys.argv) < 3:
        print("Usage: make_pptx.py <input.json> <output.pptx> [publish_date YYYY.MM.DD]", file=sys.stderr)
        sys.exit(1)

    json_path, out_path = sys.argv[1], sys.argv[2]
    with open(json_path, encoding="utf-8") as f:
        data = json.load(f)

    ws = data["week_start"].replace("-", ".")
    we = data["week_end"].replace("-", ".")
    date_range = f"{ws} ~ {we}"
    publish_date = sys.argv[3] if len(sys.argv) > 3 else we
    summary = data.get("summary", [])
    retail_items = data.get("retail", [])
    other_items = data.get("other", [])

    prs = Presentation()
    prs.slide_width = Inches(W)

    # ---- pre-compute layout heights on a throwaway slide-size guess, then fix ----
    header_h = 1.55
    orange_h = 0.05
    summary_pad = 0.28
    summary_bullet_w = W - 1.6
    summary_line_hs = [lines_needed(f"●  {line}", summary_bullet_w, 12.5, bold=True) * 0.26 for line in summary]
    summary_box_h = summary_pad * 2 + 0.30 + sum(summary_line_hs) + max(0, len(summary) - 1) * 0.04

    # rough per-item height estimate to size section boxes before drawing
    def section_height_estimate(items):
        h = 0.46 + 0.22 * 2 - 0.14
        for it in items:
            has_badge = bool(it.get("kurly_related"))
            headline_w = (W - 1.0) - 0.44 - 0.22 - (BADGE_W + BADGE_GAP if has_badge else 0)
            summary_w = (W - 1.0) - 0.44 - 0.22
            hl = lines_needed(it["headline"], headline_w, HEADLINE_SIZE, bold=True) * HEADLINE_LINE_H + 0.05
            sm = lines_needed(it["summary"], summary_w, SUMMARY_SIZE) * SUMMARY_LINE_H + 0.04
            src = 0.18 + 0.16
            h += hl + sm + src
        return h

    retail_h = section_height_estimate(retail_items)
    other_h = section_height_estimate(other_items)
    footer_h = 0.5
    gap = 0.22

    body_top = header_h + orange_h + 0.22
    total_h = body_top + summary_box_h + gap + retail_h + gap + other_h + gap + footer_h + 0.3
    prs.slide_height = Inches(total_h)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # background
    add_rect(slide, 0, 0, W, total_h, PAPER)

    # header
    add_gradient_rect(slide, 0, 0, W, header_h, PURPLE, PURPLE_2)
    add_rect(slide, 0, header_h, W, orange_h, ORANGE)
    add_text(slide, 0.5, 0.32, 6.5, 0.28, "WEEKLY HR TREND REPORT", 13, ORANGE, bold=True)
    add_text(slide, 0.5, 0.62, 6.5, 0.5, "노사동향 주간 리포트", 26, WHITE, bold=True)
    add_text(slide, 0.5, 1.12, 6.5, 0.32, date_range, 15, RGBColor(0xE3, 0xC6, 0xF0))

    cursor = body_top

    # summary box
    add_rect(slide, 0.5, cursor, W - 1.0, summary_box_h, RETAIL_BG, rounded=True, radius=0.06)
    add_rect(slide, 0.5, cursor, 0.06, summary_box_h, PURPLE)
    add_text(slide, 0.72, cursor + summary_pad - 0.06, W - 1.44, 0.26, "📌 이번 주 한눈에 보기", 11, PURPLE, bold=True)
    line_y = cursor + summary_pad + 0.24
    for line, line_h in zip(summary, summary_line_hs):
        add_text(slide, 0.9, line_y, summary_bullet_w, line_h, f"●  {line}", 12.5, INK, bold=True)
        line_y += line_h + 0.04
    cursor += summary_box_h + gap

    # retail section
    add_rect(slide, 0.5, cursor, W - 1.0, retail_h, WHITE, line=LINE, rounded=True, radius=0.05)
    build_section(slide, 0.5, cursor, W - 1.0, "🏬 유통업계", "RETAIL", retail_items, PURPLE, RETAIL_BG, PURPLE)
    cursor += retail_h + gap

    # other section
    add_rect(slide, 0.5, cursor, W - 1.0, other_h, WHITE, line=LINE, rounded=True, radius=0.05)
    build_section(slide, 0.5, cursor, W - 1.0, "🏢 타 산업군", "OTHER", other_items, ORANGE_TEXT, OTHER_BG, ORANGE)
    cursor += other_h + gap

    # footer
    add_text(slide, 0.5, cursor, 3.5, 0.24, "노사동향 위클리", 10, SOURCE_GRAY)
    add_text(slide, W - 4.0, cursor, 3.5, 0.24, f"발행일 {publish_date}", 10, SOURCE_GRAY, align=PP_ALIGN.RIGHT)

    # ---- second slide: 담당자 코멘트 (same page size) ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide2, 0, 0, W, total_h, PAPER)
    box_margin = 0.5
    box_w = W - box_margin * 2
    box_h = total_h - box_margin * 2
    add_rect(slide2, box_margin, box_margin, box_w, box_h, WHITE, line=DASH_BORDER, rounded=True, radius=0.04)

    avatar_d = 0.55
    avatar_path = Path(out_path).resolve().parent / "assets" / "avatar.jpg"
    if avatar_path.exists():
        add_circular_picture(slide2, avatar_path, box_margin + 0.35, box_margin + 0.35, avatar_d)
    else:
        av = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(box_margin + 0.35), Inches(box_margin + 0.35), Inches(avatar_d), Inches(avatar_d))
        av.fill.solid()
        av.fill.fore_color.rgb = AVATAR_BG
        av.line.fill.background()
        no_shadow(av)
        av_tf = av.text_frame
        av_tf.margin_left = 0; av_tf.margin_right = 0; av_tf.margin_top = 0; av_tf.margin_bottom = 0
        av_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        av_p = av_tf.paragraphs[0]
        av_p.alignment = PP_ALIGN.CENTER
        av_run = av_p.add_run()
        av_run.text = "🧑‍💼"
        av_run.font.size = Pt(22)

    add_text(slide2, box_margin + 0.35 + avatar_d + 0.18, box_margin + 0.35, 4.5, 0.4,
              "📝 담당자 코멘트", 18, INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    content_top = box_margin + 0.35 + avatar_d + 0.3
    add_rect(slide2, box_margin + 0.3, content_top, box_w - 0.6, 0.01, LINE)
    comment_box = slide2.shapes.add_textbox(Inches(box_margin + 0.3), Inches(content_top + 0.2),
                                              Inches(box_w - 0.6), Inches(box_h - (content_top - box_margin) - 0.5))
    ctf = comment_box.text_frame
    ctf.word_wrap = True
    p = ctf.paragraphs[0]
    run = p.add_run()
    run.text = "여기에 의견을 작성하세요."
    set_font(run, 14, RGBColor(0xA7, 0xAE, 0xB8), italic=True)

    prs.save(out_path)
    print("Saved", out_path)


if __name__ == "__main__":
    main()
